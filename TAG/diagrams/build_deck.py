"""Build the BCPC Fabric TAG deck from the approved template.

Preserves the template's narrative and branding: slides 1-3, 7 (approvals) and 8 (appendix)
are left untouched; slides 4-6 are repurposed; diagram and table slides are cloned from
existing slides so logos, footer and page numbering come along.
"""
import copy, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
OUT_DIR = os.path.join(HERE, "out")
TAG = r"C:\Users\Shuo\OneDrive\文档\bcpensionshackathon\TAG"
SRC = os.path.join(TAG, "Fabric POC - TAG-compressed.pptx")
DST = os.path.join(TAG, "Fabric POC - TAG - Platform Direction (DRAFT).pptx")

TEAL     = RGBColor(0x00, 0x77, 0x8A)
TEAL_DK  = RGBColor(0x28, 0x52, 0x5B)
GREY     = RGBColor(0x4D, 0x4D, 0x4F)
GREY_MID = RGBColor(0x7A, 0x7C, 0x80)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
TINT     = RGBColor(0xE8, 0xF2, 0xF4)
TINT2    = RGBColor(0xF4, 0xF7, 0xF8)
RED      = RGBColor(0xA0, 0x3A, 0x3A)
FONT     = "Arial"

prs = Presentation(SRC)
SW, SH = prs.slide_width, prs.slide_height


# ----------------------------------------------------------------- utilities
def clone_slide(prs, index):
    """Deep-copy a slide (shapes + relationships) and append it to the deck."""
    src = prs.slides[index]
    dst = prs.slides.add_slide(src.slide_layout)
    for shp in list(dst.shapes):                       # drop layout placeholders
        shp._element.getparent().remove(shp._element)
    # map source rIds -> new rIds on the destination part
    rid_map = {}
    for rid, rel in src.part.rels.items():
        if rel.is_external:
            new = dst.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            new = dst.part.relate_to(rel.target_part, rel.reltype)
        rid_map[rid] = new
    for shp in src.shapes:
        el = copy.deepcopy(shp._element)
        for attr in el.iter():
            for k, v in list(attr.attrib.items()):
                if k.endswith("}id") or k.endswith("}embed") or k.endswith("}link"):
                    if v in rid_map:
                        attr.set(k, rid_map[v])
        dst.shapes._spTree.append(el)
    return dst


def reorder(prs, order):
    """order = list of current slide indexes, in the sequence wanted."""
    lst = prs.slides._sldIdLst
    ids = list(lst)
    for e in ids:
        lst.remove(e)
    for i in order:
        lst.append(ids[i])


def shape_by_name(slide, name):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    return None


def drop(slide, *names):
    for n in names:
        sh = shape_by_name(slide, n)
        if sh is not None:
            sh._element.getparent().remove(sh._element)


def spaced(s):
    return " ".join(s.upper())


def set_title(slide, name, line1, line2=None):
    """Rewrite the title placeholder, keeping the template's run formatting."""
    sh = shape_by_name(slide, name)
    tf = sh.text_frame
    paras = tf.paragraphs
    def write(par, text):
        runs = par.runs
        if not runs:
            r = par.add_run(); r.text = text; return
        runs[0].text = text
        for r in runs[1:]:
            r._r.getparent().remove(r._r)
    write(paras[0], line1)
    if line2 is not None and len(paras) > 1:
        write(paras[1], spaced(line2))
    elif line2 is not None:
        par = tf.add_paragraph(); r = par.add_run(); r.text = spaced(line2)
        r.font.size = Pt(11); r.font.name = FONT; r.font.color.rgb = TEAL
    return sh


def add_pic(slide, png, top_in, width_in):
    from PIL import Image
    w, h = Image.open(png).size
    W = Inches(width_in)
    H = int(W * h / w)
    left = int((SW - W) / 2)
    return slide.shapes.add_picture(png, left, Inches(top_in), width=W, height=H)


def style_cell(cell, text, *, bold=False, size=11, color=GREY, fill=None, align=PP_ALIGN.LEFT,
               italic=False):
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        par = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        par.alignment = align
        par.space_after = Pt(2)
        run = par.add_run(); run.text = ln
        f = run.font
        f.name = FONT; f.size = Pt(size); f.bold = bold; f.italic = italic
        f.color.rgb = color
    cell.vertical_anchor = MSO_ANCHOR.TOP
    cell.margin_left = Inches(0.10); cell.margin_right = Inches(0.10)
    cell.margin_top = Inches(0.06); cell.margin_bottom = Inches(0.06)
    if fill is not None:
        cell.fill.solid(); cell.fill.fore_color.rgb = fill
    else:
        cell.fill.background()


def add_table(slide, rows, cols, left_in, top_in, width_in, height_in, col_widths=None):
    gf = slide.shapes.add_table(rows, cols, Inches(left_in), Inches(top_in),
                                Inches(width_in), Inches(height_in))
    tbl = gf.table
    tbl.first_row = False          # suppress the banded built-in style
    tbl.horz_banding = False
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Emu(int(Inches(width_in) * cw / total))
    return tbl


def note(slide, text, top_in, size=11, color=GREY_MID, left_in=0.79, width_in=11.8, bold=False):
    tb = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(0.4))
    tf = tb.text_frame; tf.word_wrap = True
    par = tf.paragraphs[0]
    run = par.add_run(); run.text = text
    run.font.name = FONT; run.font.size = Pt(size); run.font.color.rgb = color
    run.font.bold = bold
    return tb


# =============================================================== SLIDE 4 -> SCOPE
s4 = prs.slides[3]
set_title(s4, "object 4", "Scope & Approach", "what the proof of concept covers")
drop(s4, "object 5")                                   # template placeholder table

hdr = ["Phase 1  ·  Platform foundation", "Phase 2  ·  Pilot data products",
       "Phase 3  ·  Evaluate & recommend"]
body = [
    ["Stand up the F64 capacity, DEV / UAT / PROD workspaces, Entra groups, service "
     "principals and Key Vaults. Establish private connectivity through the VNet and "
     "on-premises data gateways.",
     "Onboard a small number of representative sources end to end — Bronze → Silver → Gold "
     "— with data quality, masking and row / column security applied from configuration.",
     "Measure against the success criteria, confirm the security and governance model with "
     "ITSO, and return to leadership with a clear recommendation."],
    ["Environments provisioned from code, isolated per environment, with no standing human "
     "access to production.",
     "Two or three governed Gold data products and a published Power BI report, built on "
     "masked or synthetic data only.",
     "A documented platform design, cost model and go / no-go recommendation."],
]
rowlab = ["Activities", "Outcome"]

t = add_table(s4, 4, 3, 0.79, 1.72, 11.81, 4.30)
for c in range(3):
    style_cell(t.cell(0, c), hdr[c], bold=True, size=12, color=WHITE, fill=TEAL,
               align=PP_ALIGN.CENTER)
for r in range(2):
    for c in range(3):
        style_cell(t.cell(r + 1, c), body[r][c], size=10.5,
                   fill=TINT2 if r == 0 else WHITE)
t.rows[0].height = Inches(0.42)
t.rows[1].height = Inches(1.55)
t.rows[2].height = Inches(1.25)
# merged footer row
t.cell(3, 0).merge(t.cell(3, 2))
style_cell(t.cell(3, 0),
           "Out of scope for this phase:  migrating production reporting off the current "
           "platform · decommissioning any existing system · loading unmasked member data "
           "into a non-production environment · any public-internet data path.",
           size=10.5, color=RED, fill=RGBColor(0xFD, 0xEC, 0xEC))
t.rows[3].height = Inches(0.55)

# ===================================================== SLIDE 5 -> ARCHITECTURE
s5 = prs.slides[4]
set_title(s5, "object 2", "High Level Architecture", "logical architecture")
add_pic(s5, os.path.join(OUT_DIR, "d1_architecture.png"), 1.62, 12.30)

# ============================================ SLIDE 6 -> RECOMMENDED SOLUTION
s6 = prs.slides[5]
drop(s6, "object 6", "object 7", "object 4")           # scoring table, badge, criteria label
set_title(s6, "object 5", "Recommended Solution & Approach")

rec = [
    ("Proceed with the Microsoft Fabric proof of concept",
     "Build on the architecture in this deck: one F64 capacity in a Canadian region, a "
     "medallion lakehouse on OneLake, and orchestration driven from configuration held in git."),
    ("Isolate first, prove second",
     "Three separated environments with their own workspaces, identities and secrets. "
     "Production data never moves downward; lower environments run on masked or synthetic data."),
    ("Security and governance in place before scale",
     "Group-based access through Entra, Conditional Access on every session, just-in-time "
     "privileged access, and row / column security applied from configuration — not by hand."),
    ("Private connectivity on every path",
     "The managed VNet data gateway as the default, with a small hardened on-premises gateway "
     "cluster only where Microsoft requires it. No public-internet route to BCPC data."),
    ("Add governance tooling only when a requirement calls for it",
     "Start with what Fabric provides. Microsoft Purview and Sentinel are treated as optional "
     "additions, evaluated against a specific compliance or SOC requirement."),
]
t6 = add_table(s6, len(rec), 2, 0.44, 1.72, 12.37, 4.40, col_widths=[4.2, 8.2])
for i, (head, detail) in enumerate(rec):
    style_cell(t6.cell(i, 0), head, bold=True, size=11.5, color=TEAL_DK,
               fill=TINT if i % 2 == 0 else TINT2)
    style_cell(t6.cell(i, 1), detail, size=10.5,
               fill=WHITE if i % 2 else RGBColor(0xFA, 0xFC, 0xFD))
    t6.rows[i].height = Inches(0.88)

# ======================================================== NEW DIAGRAM SLIDES
diagram_slides = [
    ("d2_environments.png", "Workspace & Environments",
     "environments and deployment"),
    ("d3_security.png", "Security, Access & Governance", "identity, access and data protection"),
    ("d4_gateway.png", "Connectivity & Gateway", "private paths from source to fabric"),
]
new_diag = []
for png, t1, t2 in diagram_slides:
    sl = clone_slide(prs, 4)                            # clone the architecture slide's chrome
    for sh in list(sl.shapes):                          # strip the cloned picture
        if sh.shape_type == 13:
            sh._element.getparent().remove(sh._element)
    set_title(sl, "object 2", t1, t2)
    add_pic(sl, os.path.join(OUT_DIR, png), 1.62, 12.30)
    new_diag.append(prs.slides.index(sl))

# ==================================== NEW TABLE SLIDES — decisions / risks
s_dec = clone_slide(prs, 3)
drop(s_dec, *[sh.name for sh in s_dec.shapes if sh.has_table])
set_title(s_dec, "object 4", "Key Decisions & Assumptions", "what we are asking tag to endorse")

dec = [
    ("Capacity & residency", "One Fabric F64 capacity hosted in a Canadian region; gateway "
     "virtual machines hosted in the same region."),
    ("Data architecture", "A medallion lakehouse on OneLake (Bronze → Silver → Gold). A "
     "warehouse is added only if a workload genuinely needs T-SQL write."),
    ("Configuration as code", "Sources, quality rules, models and security policies are held "
     "as configuration in git and promoted between environments — not hand-built per source."),
    ("Environment isolation", "Separate workspaces, Entra groups, service principals, key "
     "vaults and source connections for DEV, UAT and PROD."),
    ("Access model", "Access granted to Entra security groups only, with Conditional Access "
     "on every session and just-in-time elevation for privileged and production roles."),
    ("Connectivity", "The managed VNet data gateway is the default; a hardened, highly "
     "available on-premises gateway is used where Microsoft requires it (Oracle mirroring)."),
    ("Governance tooling", "Begin with Fabric's built-in governance. Purview and Sentinel are "
     "optional, added against a specific requirement."),
]
t7 = add_table(s_dec, len(dec), 2, 0.79, 1.62, 11.81, 3.55, col_widths=[3.1, 8.7])
for i, (k, v) in enumerate(dec):
    style_cell(t7.cell(i, 0), k, bold=True, size=10.5, color=TEAL_DK,
               fill=TINT if i % 2 == 0 else TINT2)
    style_cell(t7.cell(i, 1), v, size=10,
               fill=WHITE if i % 2 else RGBColor(0xFA, 0xFC, 0xFD))
    t7.rows[i].height = Inches(0.44)

note(s_dec, "Assumptions", 5.32, size=12, color=TEAL, bold=True)
note(s_dec,
     "The F64 capacity is available and assignable  ·  the existing VNet data gateway can be "
     "reused  ·  BCPC identity services can provision the required Entra groups and Conditional "
     "Access policies  ·  service principals can be permitted to use Fabric APIs, scoped to a "
     "named group  ·  the Oracle source can meet Microsoft's prerequisites for mirroring  ·  "
     "non-production environments will be supplied with masked or synthetic data  ·  a Canadian "
     "region satisfies BCPC's data-residency position.",
     5.62, size=10.5, color=GREY)

s_risk = clone_slide(prs, 3)
drop(s_risk, *[sh.name for sh in s_risk.shapes if sh.has_table])
set_title(s_risk, "object 4", "Risks & Mitigations", "how we are managing the exposure")

risks = [
    ("Sensitive member data in a new platform", "High",
     "The pilot runs in an isolated environment on masked or synthetic data. The security and "
     "access model is reviewed with ITSO before any real data is loaded. PIA and STRA are "
     "confirmed as part of this submission.", "ITE / ITSO"),
    ("Privileged access to production", "High",
     "No standing human change access to production. Changes flow only through the deployment "
     "pipeline executed by a service principal; administrator roles are elevated just in time "
     "with approval.", "Platform Owner"),
    ("Broad database grants required for Oracle mirroring", "Medium",
     "The on-premises gateway is dedicated, hardened and clustered, its credentials isolated in "
     "a per-environment key vault, and read-only least-privilege accounts used everywhere else.",
     "Platform Owner"),
    ("Capacity contention and cost overrun", "Medium",
     "Production runs on dedicated capacity with surge protection and throttling guardrails; "
     "consumption is monitored against the capacity allowance throughout the pilot.",
     "Capacity Admin"),
    ("Ungoverned workspace and report sprawl", "Medium",
     "Workspace creation is restricted to a request-and-approve process; every workspace is "
     "provisioned to standard from code, and only reviewed Gold products are certified.",
     "Data Governance"),
    ("Skills and adoption", "Low",
     "Adding a source is a configuration change rather than new code, which keeps the effort "
     "per source small; runbooks and working guides are maintained alongside the platform.",
     "Project team"),
]
t8 = add_table(s_risk, len(risks) + 1, 4, 0.79, 1.72, 11.81, 4.60,
               col_widths=[3.0, 0.95, 6.5, 1.5])
for c, h in enumerate(["Risk", "Rating", "Mitigation", "Owner"]):
    style_cell(t8.cell(0, c), h, bold=True, size=11, color=WHITE, fill=TEAL,
               align=PP_ALIGN.CENTER if c in (1, 3) else PP_ALIGN.LEFT)
t8.rows[0].height = Inches(0.34)
RATE = {"High": RED, "Medium": RGBColor(0xB2, 0x6B, 0x00), "Low": RGBColor(0x0E, 0x7C, 0x5F)}
for i, (r, rate, m, o) in enumerate(risks, start=1):
    fill = TINT2 if i % 2 else WHITE
    style_cell(t8.cell(i, 0), r, bold=True, size=10, color=TEAL_DK, fill=fill)
    style_cell(t8.cell(i, 1), rate, bold=True, size=10, color=RATE[rate], fill=fill,
               align=PP_ALIGN.CENTER)
    style_cell(t8.cell(i, 2), m, size=10, fill=fill)
    style_cell(t8.cell(i, 3), o, size=10, color=GREY_MID, fill=fill, align=PP_ALIGN.CENTER)
    t8.rows[i].height = Inches(0.71)

# ============================================================= CONTENTS SLIDE
s1 = prs.slides[0]
for name, text in [("object 9",  "Business Problem"),
                   ("object 10", "Business Goals"),
                   ("object 11", "Scope & Approach"),
                   ("object 12", "Architecture & Platform Design"),
                   ("object 13", "Recommended Solution & Approach"),
                   ("object 14", "Outcome")]:
    sh = shape_by_name(s1, name)
    if sh is None:
        continue
    # each label is ONE paragraph split across several runs — write it all to the first run
    par = sh.text_frame.paragraphs[0]
    runs = par.runs
    if not runs:
        continue
    runs[0].text = text
    for r in runs[1:]:
        r._r.getparent().remove(r._r)

# ==================================================================== REORDER
# current: 0 contents,1 problem,2 goals,3 scope,4 arch,5 recommended,6 outcome,7 appendix,
#          8,9,10 diagrams, 11 decisions, 12 risks
reorder(prs, [0, 1, 2, 3, 4, 8, 9, 10, 11, 12, 5, 6, 7])

prs.save(DST)
print("saved:", DST)
print("slides:", len(prs.slides.__iter__.__self__._sldIdLst))
